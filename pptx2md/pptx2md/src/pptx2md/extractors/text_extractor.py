from __future__ import annotations

from pptx.shapes.base import BaseShape

from pptx2md.models.block import BlockModel
from pptx2md.normalizers.text_cleaner import normalize_ppt_text


class TextExtractor:
    """
    Extract text blocks from a shape if the shape contains a text frame.

    V1 scope:
    - only text-bearing shapes
    - keep paragraph boundaries
    - keep paragraph levels
    - clean control characters from PPT text
    """

    def extract(self, shape: BaseShape, block_index: int) -> BlockModel | None:
        if not getattr(shape, "has_text_frame", False):
            return None

        text_frame = shape.text_frame
        if text_frame is None:
            return None

        paragraph_lines: list[str] = []

        for paragraph in text_frame.paragraphs:
            cleaned_text = normalize_ppt_text(paragraph.text)
            if not cleaned_text:
                continue

            level = getattr(paragraph, "level", 0)
            indent = "  " * level

            # preserve multi-line meaning inside a paragraph
            sub_lines = cleaned_text.split("\n")
            for i, sub_line in enumerate(sub_lines):
                if i == 0:
                    paragraph_lines.append(f"{indent}- {sub_line}")
                else:
                    paragraph_lines.append(f"{indent}  {sub_line}")

        if not paragraph_lines:
            return None

        full_text = "\n".join(paragraph_lines)

        placeholder_type = None
        if getattr(shape, "is_placeholder", False):
            try:
                placeholder_type = str(shape.placeholder_format.type)
            except Exception:
                placeholder_type = None

        block = BlockModel(
            block_id=f"block_{block_index}",
            block_type="text",
            role_hint=None,
            shape_id=getattr(shape, "shape_id", None),
            shape_name=getattr(shape, "name", None),
            shape_type=str(getattr(shape, "shape_type", None)),
            placeholder_type=placeholder_type,
            left=getattr(shape, "left", None),
            top=getattr(shape, "top", None),
            width=getattr(shape, "width", None),
            height=getattr(shape, "height", None),
            z_order=block_index,
            text=full_text,
            extra={}
        )

        return block