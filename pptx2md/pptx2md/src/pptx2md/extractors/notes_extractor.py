from __future__ import annotations

from pptx.slide import Slide

from pptx2md.normalizers.text_cleaner import normalize_ppt_text


class NotesExtractor:
    """
    Extract speaker notes text from a slide.
    """

    def extract(self, slide: Slide) -> str | None:
        try:
            notes_slide = slide.notes_slide
        except Exception:
            return None

        try:
            text_frame = notes_slide.notes_text_frame
        except Exception:
            return None

        if text_frame is None:
            return None

        lines: list[str] = []

        for paragraph in text_frame.paragraphs:
            cleaned = normalize_ppt_text(paragraph.text)
            if not cleaned:
                continue
            lines.append(cleaned)

        if not lines:
            return None

        return "\n".join(lines)