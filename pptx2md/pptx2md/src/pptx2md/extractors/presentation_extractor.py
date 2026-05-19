from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from pptx2md.models.document import DocumentModel
from pptx2md.extractors.slide_extractor import SlideExtractor


class PresentationExtractor:
    """
    Presentation extractor with assets directory and pptx path support.
    """

    def __init__(self, assets_dir: Path, pptx_path: Path) -> None:
        self.slide_extractor = SlideExtractor(assets_dir=assets_dir, pptx_path=pptx_path)

    def extract(self, input_file: Path) -> DocumentModel:
        prs = Presentation(str(input_file))

        document = DocumentModel(
            document_id="",
            source_path=input_file,
            source_name=input_file.name,
            title=None,
            slide_width=prs.slide_width,
            slide_height=prs.slide_height,
            slides=[],
        )

        for idx, slide in enumerate(prs.slides, start=1):
            slide_model = self.slide_extractor.extract(slide, idx)
            document.slides.append(slide_model)

        return document